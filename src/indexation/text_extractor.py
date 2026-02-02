"""
TextExtractor - Extract text content from various document formats.

This module provides text extraction capabilities for document indexing.
Supports PDF, DOCX, TXT, MD, JSON, TOML, and source code files.

Responsibilities:
- Extract raw text from documents without OCR
- Detect document language
- Calculate metadata (word count, page count, etc.)
- Handle encoding detection for text files
"""

from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional, Dict, Any, Type
import json

from src.core.logger import get_logger
from src.core.config import ConfigManager

# Lazy imports for optional dependencies
_pypdf = None
_docx = None
_pptx = None
_openpyxl = None
_odfpy = None
_pillow = None
_langdetect = None


def _get_pypdf():
    """Lazy load pypdf."""
    global _pypdf
    if _pypdf is None:
        import pypdf
        _pypdf = pypdf
    return _pypdf


def _get_docx():
    """Lazy load python-docx."""
    global _docx
    if _docx is None:
        import docx
        _docx = docx
    return _docx


def _get_pptx():
    """Lazy load python-pptx."""
    global _pptx
    if _pptx is None:
        import pptx
        _pptx = pptx
    return _pptx


def _get_openpyxl():
    """Lazy load openpyxl."""
    global _openpyxl
    if _openpyxl is None:
        import openpyxl
        _openpyxl = openpyxl
    return _openpyxl


def _get_odfpy():
    """Lazy load odfpy."""
    global _odfpy
    if _odfpy is None:
        from odf import text as odf_text
        from odf import opendocument
        _odfpy = {"text": odf_text, "opendocument": opendocument}
    return _odfpy


def _get_pillow():
    """Lazy load Pillow for EXIF extraction."""
    global _pillow
    if _pillow is None:
        from PIL import Image
        from PIL.ExifTags import TAGS, GPSTAGS
        _pillow = {"Image": Image, "TAGS": TAGS, "GPSTAGS": GPSTAGS}
    return _pillow


def _get_langdetect():
    """Lazy load langdetect."""
    global _langdetect
    if _langdetect is None:
        import langdetect
        _langdetect = langdetect
    return _langdetect


logger = get_logger(__name__)

# Default max file size in MB (can be overridden via config)
DEFAULT_MAX_FILE_SIZE_MB = 50

@dataclass
class ExtractionResult:
    """Result of text extraction from a document."""
    
    text: str
    """Extracted text content."""
    
    metadata: Dict[str, Any] = field(default_factory=dict)
    """Document metadata (pages, word_count, language, etc.)."""
    
    success: bool = True
    """Whether extraction was successful."""
    
    error: Optional[str] = None
    """Error message if extraction failed."""
    
    @property
    def word_count(self) -> int:
        """Return word count from metadata."""
        return self.metadata.get("word_count", 0)
    
    @property
    def language(self) -> Optional[str]:
        """Return detected language from metadata."""
        return self.metadata.get("language")
    
    @property
    def pages(self) -> Optional[int]:
        """Return page count from metadata."""
        return self.metadata.get("pages")


class BaseExtractor(ABC):
    """Abstract base class for text extractors."""
    
    # File extensions this extractor handles
    SUPPORTED_EXTENSIONS: set = set()
    
    @abstractmethod
    def extract(self, file_path: Path) -> ExtractionResult:
        """
        Extract text from a file.
        
        Args:
            file_path: Path to the file to extract text from.
            
        Returns:
            ExtractionResult with text and metadata.
        """
        pass
    
    @classmethod
    def can_handle(cls, file_path: Path) -> bool:
        """Check if this extractor can handle the given file."""
        return file_path.suffix.lower() in cls.SUPPORTED_EXTENSIONS


class PlainTextExtractor(BaseExtractor):
    """Extractor for plain text files (TXT, MD, LOG, etc.)."""
    
    SUPPORTED_EXTENSIONS = {".txt", ".md", ".markdown", ".rst", ".tex", ".log"}
    
    # Common encodings to try
    ENCODINGS = ["utf-8", "utf-16", "latin-1", "cp1252", "iso-8859-1"]
    
    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from plain text file."""
        text = None
        encoding_used = None
        
        # Try different encodings
        for encoding in self.ENCODINGS:
            try:
                text = file_path.read_text(encoding=encoding)
                encoding_used = encoding
                break
            except (UnicodeDecodeError, UnicodeError):
                continue
        
        if text is None:
            return ExtractionResult(
                text="",
                success=False,
                error=f"Could not decode file with any encoding: {self.ENCODINGS}"
            )
        
        # Calculate metadata
        word_count = len(text.split())
        line_count = text.count("\n") + 1
        
        # Detect language
        language = self._detect_language(text)
        
        return ExtractionResult(
            text=text,
            metadata={
                "word_count": word_count,
                "line_count": line_count,
                "language": language,
                "encoding": encoding_used,
                "file_type": file_path.suffix.lower().lstrip("."),
            }
        )
    
    def _detect_language(self, text: str) -> Optional[str]:
        """Detect language of text."""
        if not text or len(text.strip()) < 20:
            return None
        
        try:
            langdetect = _get_langdetect()
            # Use first 5000 chars for detection
            sample = text[:5000]
            return langdetect.detect(sample)
        except Exception:
            return None


class CodeExtractor(BaseExtractor):
    """Extractor for source code and data files."""
    
    SUPPORTED_EXTENSIONS = {
        # Python
        ".py", ".pyi", ".pyx",
        # JavaScript/TypeScript
        ".js", ".jsx", ".ts", ".tsx", ".mjs", ".cjs",
        # Web
        ".html", ".htm", ".css", ".scss", ".sass", ".less",
        # Data formats
        ".json", ".yaml", ".yml", ".toml", ".xml", ".csv",
        # Shell
        ".sh", ".bash", ".zsh", ".fish",
        # Other
        ".sql", ".r", ".rb", ".go", ".rs", ".c", ".cpp", ".h", ".hpp",
        ".java", ".kt", ".swift", ".m", ".mm",
    }
    
    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from source code file."""
        try:
            text = file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            try:
                text = file_path.read_text(encoding="latin-1")
            except Exception as e:
                return ExtractionResult(
                    text="",
                    success=False,
                    error=f"Failed to read file: {e}"
                )
        
        # Calculate metadata
        word_count = len(text.split())
        line_count = text.count("\n") + 1
        
        # Detect language (programming vs natural)
        lang = self._detect_language(text)
        
        return ExtractionResult(
            text=text,
            metadata={
                "word_count": word_count,
                "line_count": line_count,
                "language": lang,
                "programming_language": file_path.suffix.lower().lstrip("."),
                "file_type": "code",
            }
        )
    
    def _detect_language(self, text: str) -> Optional[str]:
        """Detect natural language in comments/strings."""
        if not text or len(text.strip()) < 50:
            return None
        
        try:
            langdetect = _get_langdetect()
            # Extract comments and strings for language detection
            sample = text[:5000]
            return langdetect.detect(sample)
        except Exception:
            return None


class JSONExtractor(BaseExtractor):
    """Extractor for JSON files with pretty formatting."""
    
    SUPPORTED_EXTENSIONS = {".json"}
    
    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract and format JSON content."""
        try:
            text = file_path.read_text(encoding="utf-8")
            # Parse and re-format for consistent output
            data = json.loads(text)
            formatted = json.dumps(data, indent=2, ensure_ascii=False)
            
            word_count = len(formatted.split())
            
            return ExtractionResult(
                text=formatted,
                metadata={
                    "word_count": word_count,
                    "file_type": "json",
                    "keys_count": self._count_keys(data),
                }
            )
        except json.JSONDecodeError as e:
            return ExtractionResult(
                text=text if 'text' in dir() else "",
                success=False,
                error=f"Invalid JSON: {e}"
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                success=False,
                error=f"Failed to read JSON: {e}"
            )
    
    def _count_keys(self, data: Any, count: int = 0) -> int:
        """Recursively count keys in JSON."""
        if isinstance(data, dict):
            count += len(data)
            for v in data.values():
                count = self._count_keys(v, count)
        elif isinstance(data, list):
            for item in data:
                count = self._count_keys(item, count)
        return count


class PDFExtractor(BaseExtractor):
    """Extractor for PDF files using pypdf."""
    
    SUPPORTED_EXTENSIONS = {".pdf"}
    
    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from PDF file."""
        try:
            pypdf = _get_pypdf()
        except ImportError:
            return ExtractionResult(
                text="",
                success=False,
                error="pypdf not installed. Run: uv pip install pypdf"
            )
        
        try:
            reader = pypdf.PdfReader(str(file_path))
            pages_text = []
            
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    pages_text.append(page_text)
            
            text = "\n\n".join(pages_text)
            word_count = len(text.split())
            
            # Detect language
            language = self._detect_language(text)
            
            # Extract PDF metadata
            pdf_metadata = {}
            if reader.metadata:
                pdf_metadata = {
                    "title": reader.metadata.get("/Title"),
                    "author": reader.metadata.get("/Author"),
                    "subject": reader.metadata.get("/Subject"),
                    "creator": reader.metadata.get("/Creator"),
                }
                # Clean None values
                pdf_metadata = {k: v for k, v in pdf_metadata.items() if v}
            
            return ExtractionResult(
                text=text,
                metadata={
                    "word_count": word_count,
                    "pages": len(reader.pages),
                    "language": language,
                    "file_type": "pdf",
                    **pdf_metadata,
                }
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                success=False,
                error=f"Failed to extract PDF: {e}"
            )
    
    def _detect_language(self, text: str) -> Optional[str]:
        """Detect language of text."""
        if not text or len(text.strip()) < 20:
            return None
        
        try:
            langdetect = _get_langdetect()
            sample = text[:5000]
            return langdetect.detect(sample)
        except Exception:
            return None


class DOCXExtractor(BaseExtractor):
    """Extractor for DOCX files using python-docx."""
    
    SUPPORTED_EXTENSIONS = {".docx"}
    
    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from DOCX file."""
        try:
            docx = _get_docx()
        except ImportError:
            return ExtractionResult(
                text="",
                success=False,
                error="python-docx not installed. Run: uv pip install python-docx"
            )
        
        try:
            doc = docx.Document(str(file_path))
            paragraphs = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        paragraphs.append(" | ".join(row_text))
            
            text = "\n\n".join(paragraphs)
            word_count = len(text.split())
            
            # Detect language
            language = self._detect_language(text)
            
            # Extract document properties
            doc_metadata = {}
            if doc.core_properties:
                props = doc.core_properties
                doc_metadata = {
                    "title": props.title,
                    "author": props.author,
                    "subject": props.subject,
                    "created": str(props.created) if props.created else None,
                    "modified": str(props.modified) if props.modified else None,
                }
                doc_metadata = {k: v for k, v in doc_metadata.items() if v}
            
            return ExtractionResult(
                text=text,
                metadata={
                    "word_count": word_count,
                    "paragraphs": len(paragraphs),
                    "tables": len(doc.tables),
                    "language": language,
                    "file_type": "docx",
                    **doc_metadata,
                }
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                success=False,
                error=f"Failed to extract DOCX: {e}"
            )
    
    def _detect_language(self, text: str) -> Optional[str]:
        """Detect language of text."""
        if not text or len(text.strip()) < 20:
            return None
        
        try:
            langdetect = _get_langdetect()
            sample = text[:5000]
            return langdetect.detect(sample)
        except Exception:
            return None


class PPTXExtractor(BaseExtractor):
    """Extractor for PowerPoint files using python-pptx."""
    
    SUPPORTED_EXTENSIONS = {".pptx"}
    
    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from PowerPoint file."""
        try:
            pptx = _get_pptx()
        except ImportError:
            return ExtractionResult(
                text="",
                success=False,
                error="python-pptx not installed. Run: uv pip install python-pptx"
            )
        
        try:
            prs = pptx.Presentation(str(file_path))
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"--- Slide {slide_num} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_text:
                                slide_content.append(" | ".join(row_text))
                if len(slide_content) > 1:
                    slides_text.append("\n".join(slide_content))
            
            text = "\n\n".join(slides_text)
            word_count = len(text.split())
            language = self._detect_language(text)
            
            return ExtractionResult(
                text=text,
                metadata={
                    "word_count": word_count,
                    "slides": len(prs.slides),
                    "language": language,
                    "file_type": "pptx",
                }
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                success=False,
                error=f"Failed to extract PPTX: {e}"
            )
    
    def _detect_language(self, text: str) -> Optional[str]:
        """Detect language of text."""
        if not text or len(text.strip()) < 20:
            return None
        try:
            langdetect = _get_langdetect()
            return langdetect.detect(text[:5000])
        except Exception:
            return None


class XLSXExtractor(BaseExtractor):
    """Extractor for Excel files using openpyxl."""
    
    SUPPORTED_EXTENSIONS = {".xlsx"}
    
    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from Excel file."""
        try:
            openpyxl = _get_openpyxl()
        except ImportError:
            return ExtractionResult(
                text="",
                success=False,
                error="openpyxl not installed. Run: uv pip install openpyxl"
            )
        
        try:
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            sheets_text = []
            total_rows = 0
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_content = [f"--- Sheet: {sheet_name} ---"]
                
                for row in sheet.iter_rows(values_only=True):
                    # Filter out empty cells and convert to strings
                    row_values = [str(cell) for cell in row if cell is not None]
                    if row_values:
                        sheet_content.append(" | ".join(row_values))
                        total_rows += 1
                
                if len(sheet_content) > 1:
                    sheets_text.append("\n".join(sheet_content))
            
            wb.close()
            
            text = "\n\n".join(sheets_text)
            word_count = len(text.split())
            language = self._detect_language(text)
            
            return ExtractionResult(
                text=text,
                metadata={
                    "word_count": word_count,
                    "sheets": len(wb.sheetnames),
                    "rows": total_rows,
                    "language": language,
                    "file_type": "xlsx",
                }
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                success=False,
                error=f"Failed to extract XLSX: {e}"
            )
    
    def _detect_language(self, text: str) -> Optional[str]:
        """Detect language of text."""
        if not text or len(text.strip()) < 20:
            return None
        try:
            langdetect = _get_langdetect()
            return langdetect.detect(text[:5000])
        except Exception:
            return None


class ODFExtractor(BaseExtractor):
    """Extractor for OpenDocument files using odfpy."""
    
    SUPPORTED_EXTENSIONS = {".odt", ".ods", ".odp"}
    
    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract text from OpenDocument file."""
        try:
            odfpy = _get_odfpy()
        except ImportError:
            return ExtractionResult(
                text="",
                success=False,
                error="odfpy not installed. Run: uv pip install odfpy"
            )
        
        try:
            doc = odfpy["opendocument"].load(str(file_path))
            
            # Extract all text elements
            paragraphs = []
            for element in doc.getElementsByType(odfpy["text"].P):
                text_content = self._get_text_recursive(element)
                if text_content.strip():
                    paragraphs.append(text_content.strip())
            
            text = "\n\n".join(paragraphs)
            word_count = len(text.split())
            language = self._detect_language(text)
            
            ext = file_path.suffix.lower()
            file_type = {".odt": "odt", ".ods": "ods", ".odp": "odp"}.get(ext, "odf")
            
            return ExtractionResult(
                text=text,
                metadata={
                    "word_count": word_count,
                    "paragraphs": len(paragraphs),
                    "language": language,
                    "file_type": file_type,
                }
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                success=False,
                error=f"Failed to extract ODF: {e}"
            )
    
    def _get_text_recursive(self, element) -> str:
        """Recursively extract text from ODF element."""
        result = []
        if hasattr(element, "data"):
            result.append(element.data)
        if hasattr(element, "childNodes"):
            for child in element.childNodes:
                result.append(self._get_text_recursive(child))
        return "".join(result)
    
    def _detect_language(self, text: str) -> Optional[str]:
        """Detect language of text."""
        if not text or len(text.strip()) < 20:
            return None
        try:
            langdetect = _get_langdetect()
            return langdetect.detect(text[:5000])
        except Exception:
            return None


class EXIFExtractor(BaseExtractor):
    """Extractor for image EXIF metadata using Pillow."""
    
    SUPPORTED_EXTENSIONS = {".jpg", ".jpeg", ".png", ".tiff", ".tif", ".webp"}
    
    def extract(self, file_path: Path) -> ExtractionResult:
        """Extract EXIF metadata from image file."""
        try:
            pillow = _get_pillow()
            Image = pillow["Image"]
            TAGS = pillow["TAGS"]
            GPSTAGS = pillow["GPSTAGS"]
        except ImportError:
            return ExtractionResult(
                text="",
                success=False,
                error="Pillow not installed. Run: uv pip install Pillow"
            )
        
        try:
            with Image.open(str(file_path)) as img:
                # Basic image info
                width, height = img.size
                format_type = img.format or file_path.suffix.upper().lstrip(".")
                
                # Extract EXIF data
                exif_data = {}
                raw_exif = img._getexif() if hasattr(img, "_getexif") else None
                
                if raw_exif:
                    for tag_id, value in raw_exif.items():
                        tag_name = TAGS.get(tag_id, str(tag_id))
                        # Handle GPS data separately
                        if tag_name == "GPSInfo":
                            gps_data = {}
                            for gps_tag_id, gps_value in value.items():
                                gps_tag_name = GPSTAGS.get(gps_tag_id, str(gps_tag_id))
                                gps_data[gps_tag_name] = self._convert_value(gps_value)
                            exif_data["gps"] = gps_data
                        else:
                            exif_data[tag_name] = self._convert_value(value)
                
                # Extract key metadata
                date_taken = exif_data.get("DateTimeOriginal") or exif_data.get("DateTime")
                camera_make = exif_data.get("Make", "")
                camera_model = exif_data.get("Model", "")
                camera = f"{camera_make} {camera_model}".strip() if camera_make or camera_model else None
                
                # Extract GPS coordinates
                gps_lat, gps_lon = self._extract_gps(exif_data.get("gps", {}))
                
                # Build searchable text from metadata
                text_parts = [f"Image: {file_path.name}"]
                text_parts.append(f"Dimensions: {width}x{height}")
                text_parts.append(f"Format: {format_type}")
                
                if date_taken:
                    text_parts.append(f"Date: {date_taken}")
                if camera:
                    text_parts.append(f"Camera: {camera}")
                if gps_lat is not None and gps_lon is not None:
                    text_parts.append(f"GPS: {gps_lat:.6f}, {gps_lon:.6f}")
                
                # Add other interesting EXIF fields
                if exif_data.get("ExposureTime"):
                    text_parts.append(f"Exposure: {exif_data['ExposureTime']}")
                if exif_data.get("FNumber"):
                    text_parts.append(f"Aperture: f/{exif_data['FNumber']}")
                if exif_data.get("ISOSpeedRatings"):
                    text_parts.append(f"ISO: {exif_data['ISOSpeedRatings']}")
                
                text = "\n".join(text_parts)
                
                metadata = {
                    "width": width,
                    "height": height,
                    "format": format_type,
                    "file_type": "image",
                    "has_exif": bool(raw_exif),
                }
                
                if date_taken:
                    metadata["date_taken"] = date_taken
                if camera:
                    metadata["camera"] = camera
                if gps_lat is not None:
                    metadata["gps_lat"] = gps_lat
                    metadata["gps_lon"] = gps_lon
                
                return ExtractionResult(text=text, metadata=metadata)
                
        except Exception as e:
            return ExtractionResult(
                text="",
                success=False,
                error=f"Failed to extract EXIF: {e}"
            )
    
    def _convert_value(self, value: Any) -> Any:
        """Convert EXIF value to JSON-serializable type."""
        if isinstance(value, bytes):
            try:
                return value.decode("utf-8", errors="ignore")
            except Exception:
                return str(value)
        elif hasattr(value, "numerator"):
            # Handle Ratio/Fraction types
            if value.denominator == 1:
                return value.numerator
            return float(value)
        elif isinstance(value, tuple):
            return [self._convert_value(v) for v in value]
        return value
    
    def _extract_gps(self, gps_data: Dict) -> tuple:
        """Extract GPS coordinates from EXIF GPS data."""
        if not gps_data:
            return None, None
        
        try:
            lat = gps_data.get("GPSLatitude")
            lat_ref = gps_data.get("GPSLatitudeRef", "N")
            lon = gps_data.get("GPSLongitude")
            lon_ref = gps_data.get("GPSLongitudeRef", "E")
            
            if lat and lon:
                lat_decimal = self._dms_to_decimal(lat)
                lon_decimal = self._dms_to_decimal(lon)
                
                if lat_ref == "S":
                    lat_decimal = -lat_decimal
                if lon_ref == "W":
                    lon_decimal = -lon_decimal
                
                return lat_decimal, lon_decimal
        except Exception:
            pass
        
        return None, None
    
    def _dms_to_decimal(self, dms: list) -> float:
        """Convert degrees/minutes/seconds to decimal degrees."""
        if len(dms) >= 3:
            degrees = float(dms[0])
            minutes = float(dms[1])
            seconds = float(dms[2])
            return degrees + minutes / 60 + seconds / 3600
        return 0.0


class TextExtractor:
    """
    Main text extractor that delegates to specialized extractors.
    
    This is the primary interface for text extraction. It automatically
    selects the appropriate extractor based on file extension.
    
    Example:
        extractor = TextExtractor()
        result = extractor.extract("/path/to/document.pdf")
        if result.success:
            print(f"Text: {result.text[:100]}...")
            print(f"Language: {result.language}")
            print(f"Word count: {result.word_count}")
    """
    
    # Registry of extractors
    EXTRACTORS: list[Type[BaseExtractor]] = [
        PDFExtractor,
        DOCXExtractor,
        PPTXExtractor,
        XLSXExtractor,
        ODFExtractor,
        EXIFExtractor,
        JSONExtractor,
        CodeExtractor,
        PlainTextExtractor,  # Must be last (fallback for .txt, .md, etc.)
    ]
    
    def __init__(self, max_file_size_mb: Optional[float] = None):
        """
        Initialize the text extractor.
        
        Args:
            max_file_size_mb: Maximum file size in MB. If None, uses config or default (50 MB).
        """
        self._extractors: Dict[str, BaseExtractor] = {}
        self._max_file_size_mb = max_file_size_mb or self._get_max_size_from_config()
        self._init_extractors()
    
    def _init_extractors(self):
        """Initialize extractor instances for each extension."""
        for extractor_class in self.EXTRACTORS:
            instance = extractor_class()
            for ext in extractor_class.SUPPORTED_EXTENSIONS:
                # First match wins (order matters)
                if ext not in self._extractors:
                    self._extractors[ext] = instance
    
    def _get_max_size_from_config(self) -> float:
        """Get max file size from config or return default."""
        try:
            config = ConfigManager("config/config.yaml")
            return config.get("indexation.max_file_size_mb", DEFAULT_MAX_FILE_SIZE_MB)
        except Exception:
            return DEFAULT_MAX_FILE_SIZE_MB
    
    def extract(self, file_path: str | Path) -> ExtractionResult:
        """
        Extract text from a file.
        
        Args:
            file_path: Path to the file to extract text from.
            
        Returns:
            ExtractionResult with text and metadata.
        """
        path = Path(file_path)
        
        # Check file exists
        if not path.exists():
            return ExtractionResult(
                text="",
                success=False,
                error=f"File not found: {path}"
            )
        
        if not path.is_file():
            return ExtractionResult(
                text="",
                success=False,
                error=f"Not a file: {path}"
            )
        
        # Check file size limit
        file_size_mb = path.stat().st_size / (1024 * 1024)
        if file_size_mb > self._max_file_size_mb:
            logger.warning(
                f"File too large: {path.name} ({file_size_mb:.1f} MB > {self._max_file_size_mb} MB limit)"
            )
            return ExtractionResult(
                text="",
                success=False,
                error=f"File too large: {file_size_mb:.1f} MB exceeds {self._max_file_size_mb} MB limit"
            )
        
        # Find appropriate extractor
        ext = path.suffix.lower()
        extractor = self._extractors.get(ext)
        
        if extractor is None:
            return ExtractionResult(
                text="",
                success=False,
                error=f"Unsupported file type: {ext}"
            )
        
        # Extract text
        try:
            result = extractor.extract(path)
            
            # Add common metadata
            result.metadata["file_path"] = str(path)
            result.metadata["file_name"] = path.name
            result.metadata["file_size"] = path.stat().st_size
            
            return result
        except Exception as e:
            logger.exception(f"Extraction failed for {path}")
            return ExtractionResult(
                text="",
                success=False,
                error=f"Extraction failed: {e}"
            )
    
    def get_supported_extensions(self) -> set:
        """Return set of all supported file extensions."""
        return set(self._extractors.keys())
    
    def can_extract(self, file_path: str | Path) -> bool:
        """Check if we can extract text from this file type."""
        path = Path(file_path)
        return path.suffix.lower() in self._extractors


# Convenience function
def extract_text(file_path: str | Path) -> ExtractionResult:
    """
    Extract text from a file (convenience function).
    
    Args:
        file_path: Path to the file.
        
    Returns:
        ExtractionResult with text and metadata.
    """
    extractor = TextExtractor()
    return extractor.extract(file_path)
